#!/usr/bin/env python3
"""
Build adams_vscode_users_conference_2026.pptx from Slidev presentation content.

Run from any directory:
    python doc/presentations/users_conference_2026/build_pptx.py

Requires: python-pptx, Pillow
    pip install python-pptx Pillow
"""

import io
import os
import re
import subprocess
import tempfile
import zipfile
from pathlib import Path

from lxml import etree
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
SCRIPT_DIR = Path(__file__).parent
ASSETS_DIR = (
    Path(__file__).parent.parent
    / "adams-vscode-intro"
    / "public"
)
POTX_PATH = Path(
    r"C:\Users\ben.thornton\Documents\Custom Office Templates"
    r"\2026 Cadence-Default Template.potx"
)
OUTPUT_PATH = SCRIPT_DIR / "adams_vscode_users_conference_2026.pptx"

# ---------------------------------------------------------------------------
# Slide dimensions (EMU) — standard 16:9
# ---------------------------------------------------------------------------
SLIDE_W = 12192000
SLIDE_H = 6858000

# ---------------------------------------------------------------------------
# Namespace constants
# ---------------------------------------------------------------------------
PPTX_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DRAW_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
REL_NS  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
P14_NS  = "http://schemas.microsoft.com/office/powerpoint/2010/main"

nsmap = {
    "p":   PPTX_NS,
    "a":   DRAW_NS,
    "r":   REL_NS,
    "p14": P14_NS,
}

# ---------------------------------------------------------------------------
# Colour palette (matching Slidev dark theme)
# ---------------------------------------------------------------------------
COL_BG      = RGBColor(0x0F, 0x0C, 0x29)   # near-black navy
COL_ACCENT  = RGBColor(0x63, 0x66, 0xF1)   # indigo
COL_TEXT    = RGBColor(0xE2, 0xE8, 0xF0)   # light grey-white
COL_MUTED   = RGBColor(0xA0, 0xAE, 0xC0)   # muted grey
COL_GREEN   = RGBColor(0x4A, 0xDE, 0x80)   # green check
COL_RED     = RGBColor(0xF8, 0x71, 0x71)   # red cross
COL_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# ---------------------------------------------------------------------------
# Helper: load .potx as Presentation by patching its content-type
# ---------------------------------------------------------------------------
def load_template(potx_path: Path) -> Presentation:
    buf = io.BytesIO()
    with zipfile.ZipFile(potx_path, "r") as z_in:
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z_out:
            for item in z_in.infolist():
                data = z_in.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(
                        b"presentationml.template",
                        b"presentationml.presentation",
                    )
                z_out.writestr(item, data)
    buf.seek(0)
    return Presentation(buf)


# ---------------------------------------------------------------------------
# Helper: get layout by name
# ---------------------------------------------------------------------------
def layout(prs: Presentation, name: str):
    for lay in prs.slide_layouts:
        if lay.name == name:
            return lay
    raise ValueError(f"Layout not found: {name!r}")


# ---------------------------------------------------------------------------
# Helper: remove all existing slides from the template
# ---------------------------------------------------------------------------
def clear_slides(prs: Presentation) -> None:
    """Remove all sample slides that come with the template."""
    xml_slides = prs.slides._sldIdLst
    slide_ids = [el.get("{%s}id" % REL_NS) for el in xml_slides]
    # Remove in reverse to avoid index shifting
    for rId in slide_ids:
        prs.part.drop_rel(rId)
    xml_slides.clear()
    # Trim the parts list so python-pptx doesn't try to write them
    prs.slides._sldIdLst.clear()


# ---------------------------------------------------------------------------
# Helper: set slide background colour
# ---------------------------------------------------------------------------
def set_bg(slide, hex_color: str = "#0f0c29") -> None:
    r, g, b = int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


# ---------------------------------------------------------------------------
# Helper: add a text box with optional bold / colour
# ---------------------------------------------------------------------------
def add_textbox(
    slide,
    text: str,
    x, y, cx, cy,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = None,
    align=PP_ALIGN.LEFT,
    font_name: str = "Calibri",
    word_wrap: bool = True,
) -> object:
    txBox = slide.shapes.add_textbox(x, y, cx, cy)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox


# ---------------------------------------------------------------------------
# Helper: generate a minimal black PNG (poster frame fallback for videos)
# ---------------------------------------------------------------------------
def _black_png_bytes(w: int = 640, h: int = 360) -> bytes:
    img = PILImage.new("RGB", (w, h), (20, 20, 20))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Helper: embed a video file into a slide
#
# Returns the shape element added to the slide's spTree.
# Approach: raw OOXML manipulation per ECMA-376 Part 1 §19.3
# ---------------------------------------------------------------------------
_shape_id_counter = [100]

def _next_id() -> int:
    _shape_id_counter[0] += 1
    return _shape_id_counter[0]


_first_frame_cache: dict = {}

def _first_frame_png_bytes(video_path: Path) -> bytes:
    """Extract the first frame of a video as PNG bytes using ffmpeg.
    Falls back to a dark grey PNG if ffmpeg fails."""
    if video_path in _first_frame_cache:
        return _first_frame_cache[video_path]
    try:
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp_path = tmp.name
        result = subprocess.run(
            ["ffmpeg", "-y", "-i", str(video_path), "-vframes", "1", "-q:v", "2", tmp_path],
            capture_output=True, timeout=15,
        )
        if result.returncode == 0:
            data = Path(tmp_path).read_bytes()
            _first_frame_cache[video_path] = data
            return data
    except Exception as exc:
        print(f"  WARNING: ffmpeg first-frame extraction failed for {video_path.name}: {exc}")
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass
    data = _black_png_bytes()
    _first_frame_cache[video_path] = data
    return data


def add_video(slide, video_path: Path, x, y, cx, cy) -> None:
    """
    Embed a video file into a slide.

    Adds:
      - The video file as a media part
      - A black poster-frame PNG as an image part
      - A <p:pic> shape that references both via relationships
      - A click-action hyperlink so PowerPoint recognises it as playable
    """
    video_path = Path(video_path)
    if not video_path.exists():
        print(f"  WARNING: video not found: {video_path}")
        return

    ext = video_path.suffix.lower()
    mime_map = {".mp4": "video/mp4", ".avi": "video/avi", ".gif": "image/gif"}
    mime = mime_map.get(ext, "video/mp4")

    video_bytes = video_path.read_bytes()
    poster_bytes = _first_frame_png_bytes(video_path)

    slide_part = slide.part

    # ---- Add video media part ------------------------------------------------
    # Two relationships to the same media file are required:
    #   rId_video  (standard OPC video rel) → used by <a:videoFile r:link>
    #   rId_media  (Microsoft media rel)    → used by <p14:media r:embed>
    video_partname = PackURI(f"/ppt/media/{video_path.name}")
    try:
        video_part = Part(video_partname, mime, slide_part.package, video_bytes)
        video_rId = slide_part.relate_to(
            video_part,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/video",
        )
        media_rId = slide_part.relate_to(
            video_part,
            "http://schemas.microsoft.com/office/2007/relationships/media",
        )
    except Exception as exc:
        print(f"  WARNING: could not add video part for {video_path.name}: {exc}")
        return

    # ---- Add poster image part -----------------------------------------------
    poster_partname = PackURI(f"/ppt/media/poster_{video_path.stem}.png")
    poster_part = Part(poster_partname, "image/png", slide_part.package, poster_bytes)
    poster_rId = slide_part.relate_to(
        poster_part,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
    )

    shape_id = _next_id()

    # ---- Build <p:pic> XML ---------------------------------------------------
    P  = PPTX_NS
    A  = DRAW_NS
    R  = REL_NS
    P14 = P14_NS

    pic_xml = (
        f'<p:pic xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{R}" xmlns:p14="{P14}">'
        f'  <p:nvPicPr>'
        f'    <p:cNvPr id="{shape_id}" name="{video_path.name}"/>'
        f'    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
        f'    <p:nvPr>'
        f'      <a:videoFile r:link="{video_rId}"/>'
        f'      <p:extLst>'
        f'        <p:ext uri="{{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}}">'
        f'          <p14:media r:embed="{media_rId}"/>'
        f'        </p:ext>'
        f'      </p:extLst>'
        f'    </p:nvPr>'
        f'  </p:nvPicPr>'
        f'  <p:blipFill>'
        f'    <a:blip r:embed="{poster_rId}"/>'
        f'    <a:stretch><a:fillRect/></a:stretch>'
        f'  </p:blipFill>'
        f'  <p:spPr>'
        f'    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>'
        f'    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'  </p:spPr>'
        f'</p:pic>'
    )

    pic_el = etree.fromstring(pic_xml)
    slide.shapes._spTree.append(pic_el)


# ---------------------------------------------------------------------------
# Helper: embed a GIF as an animated picture
#
# python-pptx's add_picture uses image/png conversion by default for GIF.
# We bypass that by adding the GIF bytes directly as image/gif so PowerPoint
# receives the raw animated GIF and plays it in presentation mode.
# ---------------------------------------------------------------------------
def add_gif(slide, gif_path: Path, x, y, cx, cy) -> None:
    gif_path = Path(gif_path)
    if not gif_path.exists():
        print(f"  WARNING: gif not found: {gif_path}")
        return

    gif_bytes = gif_path.read_bytes()
    slide_part = slide.part

    gif_partname = PackURI(f"/ppt/media/{gif_path.name}")
    gif_part = Part(gif_partname, "image/gif", slide_part.package, gif_bytes)
    img_rId = slide_part.relate_to(
        gif_part,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
    )

    shape_id = _next_id()
    P = PPTX_NS
    A = DRAW_NS
    R = REL_NS

    pic_xml = (
        f'<p:pic xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{R}">'
        f'  <p:nvPicPr>'
        f'    <p:cNvPr id="{shape_id}" name="{gif_path.name}"/>'
        f'    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
        f'    <p:nvPr/>'
        f'  </p:nvPicPr>'
        f'  <p:blipFill>'
        f'    <a:blip r:embed="{img_rId}"/>'
        f'    <a:stretch><a:fillRect/></a:stretch>'
        f'  </p:blipFill>'
        f'  <p:spPr>'
        f'    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>'
        f'    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'  </p:spPr>'
        f'</p:pic>'
    )

    pic_el = etree.fromstring(pic_xml)
    slide.shapes._spTree.append(pic_el)


# ---------------------------------------------------------------------------
# Helper: add an image via Pillow (normal static images)
# ---------------------------------------------------------------------------
def add_image(slide, img_path: Path, x, y, cx, cy) -> None:
    img_path = Path(img_path)
    if not img_path.exists():
        print(f"  WARNING: image not found: {img_path}")
        return
    slide.shapes.add_picture(str(img_path), x, y, cx, cy)


# ---------------------------------------------------------------------------
# Helper: add entrance "Appear" animations to a list of shape IDs on a slide
#
# Each shape_id appears on a separate mouse click (sequential).
# For auto-staggered (all on one click, delay-based), pass stagger_ms > 0.
# ---------------------------------------------------------------------------
def add_appear_animations(
    slide, shape_ids: list, stagger_ms: int = 0
) -> None:
    """
    Inject <p:timing> XML to make each shape appear on click (or staggered).

    shape_ids: list of integer shape IDs (sp.shape_id values)
    stagger_ms: if > 0, all shapes appear on one click with this delay between them
    """
    if not shape_ids:
        return

    P = PPTX_NS
    A = DRAW_NS

    # Build the par elements for each shape
    par_elements = []
    for i, spid in enumerate(shape_ids):
        delay = i * stagger_ms if stagger_ms > 0 else 0
        # Each shape starts on click (indefinite delay) except stagger mode
        start_cond = (
            f'<p:cond delay="{delay}"/>'
            if stagger_ms > 0
            else '<p:cond evt="onBegin" delay="indefinite"/>'
        )

        node_type = "withEffect" if stagger_ms > 0 else "clickEffect"

        ctn_id_outer = _next_id()
        ctn_id_mid   = _next_id()
        ctn_id_inner = _next_id()
        ctn_id_set   = _next_id()

        par_elements.append(
            f'<p:par xmlns:p="{P}" xmlns:a="{A}">'
            f'  <p:cTn id="{ctn_id_outer}" fill="hold" nodeType="{node_type}">'
            f'    <p:stCondLst>{start_cond}</p:stCondLst>'
            f'    <p:childTnLst>'
            f'      <p:par>'
            f'        <p:cTn id="{ctn_id_mid}" fill="hold">'
            f'          <p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'          <p:childTnLst>'
            f'            <p:par>'
            f'              <p:cTn id="{ctn_id_inner}" presetID="1" presetClass="entr"'
            f'                     presetSubtype="0" fill="hold" grpId="{i}"'
            f'                     nodeType="{node_type}">'
            f'                <p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'                <p:childTnLst>'
            f'                  <p:set>'
            f'                    <p:cBhvr>'
            f'                      <p:cTn id="{ctn_id_set}" dur="1" fill="hold"/>'
            f'                      <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            f'                      <p:attrNameLst>'
            f'                        <p:attrName>style.visibility</p:attrName>'
            f'                      </p:attrNameLst>'
            f'                    </p:cBhvr>'
            f'                    <p:to><p:strVal val="visible"/></p:to>'
            f'                  </p:set>'
            f'                </p:childTnLst>'
            f'              </p:cTn>'
            f'            </p:par>'
            f'          </p:childTnLst>'
            f'        </p:cTn>'
            f'      </p:par>'
            f'    </p:childTnLst>'
            f'  </p:cTn>'
            f'</p:par>'
        )

    # Build bldLst entries (no build attribute = appear as whole unit)
    bld_entries = "".join(
        f'<p:bldP xmlns:p="{P}" spid="{spid}" grpId="{i}" uiExpand="1"/>'
        for i, spid in enumerate(shape_ids)
    )

    root_id   = _next_id()
    seq_id    = _next_id()

    children_xml = "\n".join(par_elements)

    timing_xml = (
        f'<p:timing xmlns:p="{P}" xmlns:a="{A}">'
        f'  <p:tnLst>'
        f'    <p:par>'
        f'      <p:cTn id="{root_id}" dur="indefinite" restart="never"'
        f'             nodeType="tmRoot">'
        f'        <p:childTnLst>'
        f'          <p:seq concurrent="1" nextAc="seek">'
        f'            <p:cTn id="{seq_id}" dur="indefinite" nodeType="mainSeq">'
        f'              <p:childTnLst>'
        f'                {children_xml}'
        f'              </p:childTnLst>'
        f'            </p:cTn>'
        f'            <p:prevCondLst>'
        f'              <p:cond evt="onBegin" delay="0"><p:tn val="0"/></p:cond>'
        f'            </p:prevCondLst>'
        f'            <p:nextCondLst>'
        f'              <p:cond evt="onBegin" delay="0"><p:tn val="0"/></p:cond>'
        f'            </p:nextCondLst>'
        f'          </p:seq>'
        f'        </p:childTnLst>'
        f'      </p:cTn>'
        f'    </p:par>'
        f'  </p:tnLst>'
        f'  <p:bldLst>{bld_entries}</p:bldLst>'
        f'</p:timing>'
    )

    timing_el = etree.fromstring(timing_xml)

    # Replace any existing timing element, or append to <p:sld> root
    sld_el = slide.element  # <p:sld>
    for old in sld_el.findall("{%s}timing" % P):
        sld_el.remove(old)
    sld_el.append(timing_el)


# ---------------------------------------------------------------------------
# Helper: hide shape initially (for animation targets)
#
# NOTE: For OOXML Appear animations, shapes do NOT need to be pre-hidden in
# the XML. PowerPoint automatically hides bldLst shapes at playback start
# and reveals them when the animation fires. This function is intentionally
# a no-op — it exists so callers are explicit about intent without risk.
# ---------------------------------------------------------------------------
def hide_shape(shape) -> None:
    """No-op: PowerPoint hides animated shapes automatically via bldLst."""
    pass


# ---------------------------------------------------------------------------
# Helper: add speaker notes to a slide
# ---------------------------------------------------------------------------
def add_notes(slide, text: str) -> None:
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# ---------------------------------------------------------------------------
# Helper: add a styled section header (thin horizontal rule + subtitle)
# ---------------------------------------------------------------------------
def add_section_sub(slide, subtitle: str, y_offset=None) -> None:
    y = y_offset or int(SLIDE_H * 0.72)
    add_textbox(
        slide, subtitle,
        x=Emu(int(SLIDE_W * 0.037)),
        y=Emu(y),
        cx=Emu(int(SLIDE_W * 0.88)),
        cy=Emu(int(SLIDE_H * 0.10)),
        font_size=18,
        color=RGBColor(0xA5, 0xB4, 0xFC),  # light indigo
        align=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Helper: add a styled "code block" text box (dark bg, mono font)
# ---------------------------------------------------------------------------
def add_code_block(
    slide, code_lines: list, x, y, cx, cy, font_size: int = 13
) -> object:
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(x, y, cx, cy)
    tf = txBox.text_frame
    tf.word_wrap = False

    # Dark background on the text box frame
    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)

    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = RGBColor(0xCD, 0xD6, 0xF4)

    return txBox


# ===========================================================================
# Slide builders
# ===========================================================================

def slide_01_title(prs):
    """Title slide: Streamlining Adams Scripting"""
    slide = prs.slides.add_slide(layout(prs, "Title Slide"))
    set_bg(slide, "#0f0c29")

    # Title placeholder (idx=0)
    ph = slide.placeholders[0]
    ph.text = "Streamlining Adams Scripting"
    tf = ph.text_frame
    tf.paragraphs[0].runs[0].font.size = Pt(36)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = COL_WHITE
    tf.paragraphs[0].runs[0].font.name = "Calibri Light"

    # Subtitle placeholder (idx=1)
    ph_sub = slide.placeholders[1]
    ph_sub.text = "The Adams VS Code Extension"
    tf_sub = ph_sub.text_frame
    tf_sub.paragraphs[0].runs[0].font.size = Pt(20)
    tf_sub.paragraphs[0].runs[0].font.color.rgb = COL_ACCENT
    tf_sub.paragraphs[0].runs[0].font.name = "Calibri"

    # Adams logo (right side)
    logo_path = ASSETS_DIR / "adams-logo.png"
    if logo_path.exists():
        slide.shapes.add_picture(
            str(logo_path),
            Emu(int(SLIDE_W * 0.72)),
            Emu(int(SLIDE_H * 0.05)),
            height=Emu(int(SLIDE_H * 0.08)),
        )

    # Author
    add_textbox(
        slide, "Ben Thornton",
        x=Emu(int(SLIDE_W * 0.037)),
        y=Emu(int(SLIDE_H * 0.88)),
        cx=Emu(int(SLIDE_W * 0.40)),
        cy=Emu(int(SLIDE_H * 0.08)),
        font_size=14,
        color=COL_MUTED,
    )

    add_notes(slide, (
        "Welcome everyone. I'm going to show you something that I think will "
        "change how you write Adams scripts. Whether you're a CMD power user or "
        "just getting started with scripting, this extension brings modern editor "
        "intelligence to Adams."
    ))
    return slide


def slide_02_agenda(prs):
    """Agenda slide with 5 staggered items."""
    slide = prs.slides.add_slide(layout(prs, "Title with Tighter Content Leading"))
    set_bg(slide, "#0f0c29")

    slide.placeholders[0].text = "Agenda"
    _style_title(slide)

    items = [
        ("01", "The Editor Gap",       "Where Adams scripting tools are today"),
        ("02", "Code Editing",         "Syntax highlighting, autocomplete, hover docs, linting, code navigation"),
        ("03", "Adams Integration",    "Run in Adams, Python debugging"),
        ("04", "What's Coming",        "Teaching AI agents to use Adams"),
        ("05", "Get Started",          "Two-minute install"),
    ]

    y_start  = Emu(1300000)
    row_h    = Emu(860000)
    shape_ids = []

    for i, (num, title, sub) in enumerate(items):
        y = y_start + i * row_h

        # Numbered label
        n_box = add_textbox(
            slide, num,
            x=Emu(480000), y=y,
            cx=Emu(500000), cy=row_h,
            font_size=28, bold=True,
            color=COL_ACCENT,
            font_name="Calibri Light",
        )
        hide_shape(n_box)
        shape_ids.append(n_box.shape_id)

        # Title + sub in one box
        txBox = slide.shapes.add_textbox(Emu(1100000), y, Emu(10200000), row_h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title
        r1.font.size = Pt(18)
        r1.font.bold = True
        r1.font.color.rgb = COL_TEXT
        r1.font.name = "Calibri"
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(13)
        r2.font.color.rgb = COL_MUTED
        r2.font.name = "Calibri"
        hide_shape(txBox)
        shape_ids.append(txBox.shape_id)

    add_appear_animations(slide, shape_ids, stagger_ms=200)

    add_notes(slide, (
        "Here's what we'll cover. We'll start with the current state — what we've "
        "all been working with. Then I'll walk through the everyday improvements: "
        "autocomplete, hover docs, linting. Then the bigger workflow changes: "
        "running code in Adams, Python debugging. A quick look at what's coming "
        "with AI. And then we'll get you installed."
    ))
    return slide


def slide_03_section_editor_gap(prs):
    """Section: The Editor Gap"""
    slide = prs.slides.add_slide(layout(prs, "Segue Slide Blue"))
    slide.placeholders[0].text = "The Editor Gap"
    slide.placeholders[1].text = "Where Adams scripting tools are today"
    add_notes(slide, "Let's start by talking about what we've all been working with.")
    return slide


def slide_04_aview_editor(prs):
    """The Adams GUI Macro Editor screenshot."""
    slide = prs.slides.add_slide(layout(prs, "Title Only"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "The Adams GUI Macro Editor"
    _style_title(slide)
    add_image(
        slide, ASSETS_DIR / "aview_macro_editor.png",
        x=Emu(int(SLIDE_W * 0.10)),
        y=Emu(int(SLIDE_H * 0.18)),
        cx=Emu(int(SLIDE_W * 0.80)),
        cy=Emu(int(SLIDE_H * 0.70)),
    )
    add_notes(slide, (
        "This is the Adams macro editor. It's a text box. And an Apply button. "
        "That's it. No highlighting, no error checking, no docs. If you misspell "
        "an argument, you find out when you run it."
    ))
    return slide


def slide_05_notepadpp(prs):
    """The Notepad++ Upgrade screenshot."""
    slide = prs.slides.add_slide(layout(prs, "Title Only"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "The Notepad++ Upgrade"
    _style_title(slide)
    add_image(
        slide, ASSETS_DIR / "notepadpp.png",
        x=Emu(int(SLIDE_W * 0.10)),
        y=Emu(int(SLIDE_H * 0.18)),
        cx=Emu(int(SLIDE_W * 0.80)),
        cy=Emu(int(SLIDE_H * 0.70)),
    )
    add_notes(slide, (
        "If you wanted better, you moved to Notepad++ with a custom syntax file. "
        "You got colors. That's it. No completions, no docs, no linting. "
        "Every syntax error still costs you a round-trip to Adams View."
    ))
    return slide


def slide_06_whats_missing(prs):
    """What's Missing? — 5 click-animated bullets + 2 conclusion boxes."""
    slide = prs.slides.add_slide(layout(prs, "Title with Tighter Content Leading"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "What's Missing?"
    _style_title(slide)

    bullets = [
        ("✗", "No error checking",   "syntax errors found at runtime, not while writing"),
        ("✗", "No documentation",    "constant tab-switching to the Adams help"),
        ("✗", "No autocomplete",     "memorize every command and argument name"),
        ("✗", "No code navigation",  "find your macros by searching directories"),
        ("✗", "No debugging",        "print statements"),
    ]

    y_start = Emu(1250000)
    row_h   = Emu(780000)
    shape_ids = []

    for i, (icon, title, sub) in enumerate(bullets):
        y = y_start + i * row_h
        txBox = slide.shapes.add_textbox(Emu(480000), y, Emu(10900000), row_h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        # Red ✗
        r_icon = p.add_run()
        r_icon.text = icon + "  "
        r_icon.font.size = Pt(18)
        r_icon.font.color.rgb = COL_RED
        r_icon.font.name = "Calibri"
        # Bold title
        r_title = p.add_run()
        r_title.text = title + " — "
        r_title.font.size = Pt(18)
        r_title.font.bold = True
        r_title.font.color.rgb = COL_TEXT
        r_title.font.name = "Calibri"
        # Muted subtitle
        r_sub = p.add_run()
        r_sub.text = sub
        r_sub.font.size = Pt(18)
        r_sub.font.color.rgb = COL_MUTED
        r_sub.font.name = "Calibri"
        hide_shape(txBox)
        shape_ids.append(txBox.shape_id)

    # Conclusion box 1
    box1 = slide.shapes.add_textbox(
        Emu(480000), Emu(5200000), Emu(10900000), Emu(500000)
    )
    tf1 = box1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "Every modern language has all of this."
    r1.font.size = Pt(18)
    r1.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    r1.font.name = "Calibri"
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
    hide_shape(box1)
    shape_ids.append(box1.shape_id)

    # Conclusion box 2
    box2 = slide.shapes.add_textbox(
        Emu(480000), Emu(5780000), Emu(10900000), Emu(500000)
    )
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Now Adams does too!"
    r2.font.size = Pt(22)
    r2.font.bold = True
    r2.font.color.rgb = RGBColor(0xA5, 0xB4, 0xFC)
    r2.font.name = "Calibri Light"
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x4A)
    hide_shape(box2)
    shape_ids.append(box2.shape_id)

    add_appear_animations(slide, shape_ids)

    add_notes(slide, (
        "Let's be honest about what's missing. No error checking — you find typos "
        "at runtime. No docs — you're constantly switching to the help browser. "
        "No autocomplete — you have to memorize command names and argument lists. "
        "No navigation — finding a macro means searching folders. And debugging? "
        "Print statements and prayer.\n\n"
        "Every single syntax error costs you a round-trip to Adams View. "
        "Write, run, fail, fix, repeat. What if your editor could catch those "
        "problems before you ever hit run?"
    ))
    return slide


def slide_07_section_code_editing(prs):
    """Section: Code Editing"""
    slide = prs.slides.add_slide(layout(prs, "Segue Slide Blue"))
    slide.placeholders[0].text = "Code Editing"
    slide.placeholders[1].text = "How the MSC Adams extension for VS Code helps you read and write code"
    add_notes(slide, (
        "Let's see what changes when your editor actually understands Adams."
    ))
    return slide


def slide_08_syntax_highlighting(prs):
    """Syntax Highlighting — two side-by-side screenshots."""
    slide = prs.slides.add_slide(layout(prs, "Title with Tighter Content Leading"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Syntax Highlighting"
    _style_title(slide)

    mid_x = int(SLIDE_W * 0.50)
    img_y  = int(SLIDE_H * 0.18)
    img_h  = int(SLIDE_H * 0.72)
    img_w  = int(SLIDE_W * 0.44)
    gap    = int(SLIDE_W * 0.03)

    # Left: dark theme
    add_textbox(slide, "Dark theme",
        x=Emu(int(SLIDE_W * 0.04)), y=Emu(img_y - 280000),
        cx=Emu(img_w), cy=Emu(280000),
        font_size=12, color=COL_MUTED, align=PP_ALIGN.CENTER)
    add_image(slide, ASSETS_DIR / "syntax_highlighting.png",
        x=Emu(int(SLIDE_W * 0.04)), y=Emu(img_y),
        cx=Emu(img_w), cy=Emu(img_h))

    # Right: light theme
    add_textbox(slide, "Light theme",
        x=Emu(mid_x + gap), y=Emu(img_y - 280000),
        cx=Emu(img_w), cy=Emu(280000),
        font_size=12, color=COL_MUTED, align=PP_ALIGN.CENTER)
    add_image(slide, ASSETS_DIR / "syntax_highlighting-light.png",
        x=Emu(mid_x + gap), y=Emu(img_y),
        cx=Emu(img_w), cy=Emu(img_h))

    add_notes(slide, (
        "VS Code with semantic token highlighting — the editor distinguishes "
        "commands, arguments, values, and even valid vs invalid names with "
        "different colors. Works in both dark and light themes."
    ))
    return slide


def slide_09_autocomplete_cmd(prs):
    """Autocomplete — Adams CMD video."""
    slide = prs.slides.add_slide(layout(prs, "Title Only"))
    set_bg(slide, "#0f0c29")
    ph = slide.placeholders[0]
    ph.text = "Autocomplete"
    _style_title(slide)
    add_textbox(slide, "Adams CMD",
        x=Emu(446313), y=Emu(730000), cx=Emu(int(SLIDE_W * 0.8)),
        cy=Emu(300000), font_size=16, color=COL_MUTED)
    add_video(slide, ASSETS_DIR / "cmd_autocomplete.mp4",
        x=Emu(int(SLIDE_W * 0.05)), y=Emu(1200000),
        cx=Emu(int(SLIDE_W * 0.90)), cy=Emu(int(SLIDE_H * 0.76)))
    add_notes(slide, (
        "Start typing a command name, and the editor shows you completions with "
        "the full argument list. Tab-complete into a template. You don't need to "
        "memorize argument names — the editor knows them."
    ))
    return slide


def slide_10_autocomplete_python(prs):
    """Autocomplete — Python video."""
    slide = prs.slides.add_slide(layout(prs, "Title Only"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Autocomplete"
    _style_title(slide)
    add_textbox(slide, "Python",
        x=Emu(446313), y=Emu(730000), cx=Emu(int(SLIDE_W * 0.8)),
        cy=Emu(300000), font_size=16, color=COL_MUTED)
    add_video(slide, ASSETS_DIR / "python_autocomplete.mp4",
        x=Emu(int(SLIDE_W * 0.05)), y=Emu(1200000),
        cx=Emu(int(SLIDE_W * 0.90)), cy=Emu(int(SLIDE_H * 0.76)))
    add_notes(slide, "Full autocomplete for the Adams Python API — every class, method, and argument.")
    return slide


def slide_11_hover_cmd(prs):
    """Hover Documentation — Adams CMD video."""
    slide = prs.slides.add_slide(layout(prs, "Title Only"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Hover Documentation"
    _style_title(slide)
    add_textbox(slide, "Adams CMD",
        x=Emu(446313), y=Emu(730000), cx=Emu(int(SLIDE_W * 0.8)),
        cy=Emu(300000), font_size=16, color=COL_MUTED)
    add_video(slide, ASSETS_DIR / "hover_cmd.mp4",
        x=Emu(int(SLIDE_W * 0.05)), y=Emu(1200000),
        cx=Emu(int(SLIDE_W * 0.90)), cy=Emu(int(SLIDE_H * 0.76)))
    add_notes(slide, (
        "Hover over any Adams function — DX, STEP, IMPACT — and you get the full "
        "documentation inline. Arguments, format, examples. Works for built-in "
        "commands, custom macros, and abbreviated forms."
    ))
    return slide


def slide_12_hover_python(prs):
    """Hover Documentation — Python video."""
    slide = prs.slides.add_slide(layout(prs, "Title Only"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Hover Documentation"
    _style_title(slide)
    add_textbox(slide, "Python",
        x=Emu(446313), y=Emu(730000), cx=Emu(int(SLIDE_W * 0.8)),
        cy=Emu(300000), font_size=16, color=COL_MUTED)
    add_video(slide, ASSETS_DIR / "hover_python.mp4",
        x=Emu(int(SLIDE_W * 0.05)), y=Emu(1200000),
        cx=Emu(int(SLIDE_W * 0.90)), cy=Emu(int(SLIDE_H * 0.76)))
    add_notes(slide, (
        "The same hover docs experience for Adams Python scripts. Hover over any "
        "Adams Python API method and you get the full docstring inline."
    ))
    return slide


def slide_13_linting_demo(prs):
    """Linting — before/after code blocks (animated)."""
    slide = prs.slides.add_slide(layout(prs, "Title with Tighter Content Leading"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Linting"
    _style_title(slide)
    add_textbox(slide, "Spellcheck for code",
        x=Emu(446313), y=Emu(730000), cx=Emu(int(SLIDE_W * 0.8)),
        cy=Emu(300000), font_size=16, color=COL_MUTED)

    code_correct = [
        "marker create  &",
        "   marker_name = .model.PART_1.cm  &",
        "   location    = 0, 0, 0  &",
        "   orientation = 0, 0, 0",
    ]
    code_typo = [
        "marker create  &",
        "   marker_name = .model.PART_1.cm  &",
        "   locaton     = 0, 0, 0  &   ← misspelled argument",
        "   orientation = 0, 0, 0",
    ]

    bx1 = add_code_block(slide, code_correct,
        x=Emu(int(SLIDE_W * 0.05)), y=Emu(1500000),
        cx=Emu(int(SLIDE_W * 0.88)), cy=Emu(1600000))

    bx2 = add_code_block(slide, code_typo,
        x=Emu(int(SLIDE_W * 0.05)), y=Emu(3400000),
        cx=Emu(int(SLIDE_W * 0.88)), cy=Emu(1600000))

    # Label boxes
    lbl1 = add_textbox(slide, "✓  Before linter — looks fine",
        x=Emu(int(SLIDE_W * 0.05)), y=Emu(1280000),
        cx=Emu(int(SLIDE_W * 0.40)), cy=Emu(260000),
        font_size=12, color=COL_GREEN)
    lbl2 = add_textbox(slide, "✗  After linter — typo flagged",
        x=Emu(int(SLIDE_W * 0.05)), y=Emu(3180000),
        cx=Emu(int(SLIDE_W * 0.40)), cy=Emu(260000),
        font_size=12, color=COL_RED)

    hide_shape(bx2)
    hide_shape(lbl2)
    add_appear_animations(slide, [bx2.shape_id, lbl2.shape_id])

    add_notes(slide, (
        "The editor uses semantic tokens to color valid and invalid argument names "
        "differently. 'locaton' is misspelled — the wavy red underline is the "
        "linter kicking in."
    ))
    return slide


def slide_14_linting_cmd(prs):
    """Linting — Adams CMD video + severity legend."""
    slide = prs.slides.add_slide(layout(prs, "Title Only"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Linting"
    _style_title(slide)
    add_textbox(slide, "Adams CMD",
        x=Emu(446313), y=Emu(730000), cx=Emu(int(SLIDE_W * 0.8)),
        cy=Emu(300000), font_size=16, color=COL_MUTED)

    # Video takes most of the slide
    add_video(slide, ASSETS_DIR / "linting.mp4",
        x=Emu(int(SLIDE_W * 0.04)), y=Emu(1200000),
        cx=Emu(int(SLIDE_W * 0.68)), cy=Emu(int(SLIDE_H * 0.76)))

    # Legend on the right
    legend = [
        (COL_RED,                    "Error   E000",  "e.g. unknown command"),
        (RGBColor(0xF5, 0xC2, 0x42), "Warning W000",  "e.g. object name omitted"),
        (RGBColor(0x6A, 0xB0, 0xF5), "Info    I000",  "e.g. hardcoded Adams ID"),
    ]
    lx = Emu(int(SLIDE_W * 0.75))
    ly = Emu(1400000)
    lcy = Emu(900000)
    for col, label, desc in legend:
        box = slide.shapes.add_textbox(lx, ly, Emu(int(SLIDE_W * 0.22)), lcy)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = col
        r.font.name = "Courier New"
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = COL_MUTED
        r2.font.name = "Calibri"
        ly += lcy + Emu(200000)

    add_notes(slide, (
        "The linter checks your code as you type. Unknown command? Red squiggle. "
        "Invalid argument name? Error immediately. You see the problem before you "
        "ever hit run. Every red squiggle is a round-trip to Adams View you just saved."
    ))
    return slide


def slide_15_linting_python(prs):
    """Linting — Python video."""
    slide = prs.slides.add_slide(layout(prs, "Title Only"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Linting"
    _style_title(slide)
    add_textbox(slide, "Python",
        x=Emu(446313), y=Emu(730000), cx=Emu(int(SLIDE_W * 0.8)),
        cy=Emu(300000), font_size=16, color=COL_MUTED)
    add_video(slide, ASSETS_DIR / "python_linting.mp4",
        x=Emu(int(SLIDE_W * 0.05)), y=Emu(1200000),
        cx=Emu(int(SLIDE_W * 0.90)), cy=Emu(int(SLIDE_H * 0.76)))
    add_notes(slide, (
        "Python scripts get the same treatment via Pylance — type errors, missing "
        "attributes, and invalid API usage surfaced as you write."
    ))
    return slide


def slide_16_code_nav(prs):
    """Code Navigation — Adams CMD video."""
    slide = prs.slides.add_slide(layout(prs, "Title Only"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Code Navigation"
    _style_title(slide)
    add_textbox(slide, "Adams CMD",
        x=Emu(446313), y=Emu(730000), cx=Emu(int(SLIDE_W * 0.8)),
        cy=Emu(300000), font_size=16, color=COL_MUTED)
    add_video(slide, ASSETS_DIR / "cmd_linked_refs.mp4",
        x=Emu(int(SLIDE_W * 0.05)), y=Emu(1200000),
        cx=Emu(int(SLIDE_W * 0.90)), cy=Emu(int(SLIDE_H * 0.76)))
    add_notes(slide, (
        "Click on any name — a part, a marker, a variable, a macro — and Go to "
        "Definition jumps straight to where it's defined. Find All References "
        "shows every place it's used across the entire workspace."
    ))
    return slide


def slide_17_custom_macros(prs):
    """Compatible with Custom Macros — video + feature chips."""
    slide = prs.slides.add_slide(layout(prs, "Title Only"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Compatible with Custom Macros"
    _style_title(slide)

    add_video(slide, ASSETS_DIR / "works_with_custom_macros.mp4",
        x=Emu(int(SLIDE_W * 0.04)), y=Emu(1100000),
        cx=Emu(int(SLIDE_W * 0.92)), cy=Emu(int(SLIDE_H * 0.65)))

    chips = ["✓ Autocomplete", "✓ Hover docs", "✓ Linting",
             "✓ Go to Definition", "✓ Find All References"]
    cx_each = Emu(int(SLIDE_W * 0.17))
    gap     = Emu(int(SLIDE_W * 0.015))
    total_w = len(chips) * cx_each + (len(chips) - 1) * gap
    x_start = Emu((SLIDE_W - int(total_w)) // 2)
    y_chip  = Emu(int(SLIDE_H * 0.875))
    for chip in chips:
        box = slide.shapes.add_textbox(x_start, y_chip, cx_each, Emu(380000))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = chip
        r.font.size = Pt(12)
        r.font.color.rgb = COL_GREEN
        r.font.name = "Calibri"
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x05, 0x2A, 0x14)
        x_start += cx_each + gap

    add_notes(slide, (
        "Everything we just showed works for your custom macros too. "
        "Your macros appear in the autocomplete list. Hover over a call and you "
        "see the help string. The linter knows they're valid commands. "
        "Go to Definition and Find All References work across your entire workspace."
    ))
    return slide


def slide_18_file_types(prs):
    """Supported File Types — feature table."""
    slide = prs.slides.add_slide(layout(prs, "Title with Tighter Content Leading"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Supported File Types"
    _style_title(slide)

    headers = ["Extension", "Syntax\nHighlighting", "Autocomplete\n& Hover",
               "Linting", "Code\nNavigation", "Run in\nAdams"]
    rows = [
        (".cmd / .mac",              "✓", "✓", "✓", "✓", "✓"),
        (".py",                      "✓", "✓", "✓", "✓", "✓"),
        (".adm / .acf",              "✓", "—", "—", "—", "—"),
        (".msg / aview.log",         "✓", "—", "—", "—", "—"),
        ("Template Files (Time Orbit)","✓", "—", "—", "—", "—"),
    ]

    table_x = Emu(446313)
    table_y = Emu(1200000)
    table_w = Emu(11255827)
    table_h = Emu(5100000)

    tbl = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        table_x, table_y, table_w, table_h
    ).table

    col_widths = [Emu(2500000)] + [Emu(int((table_w - 2500000) / (len(headers) - 1)))] * (len(headers) - 1)
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    def _cell(r, c, text, bold=False, color=None, center=False):
        cell = tbl.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = (
            RGBColor(0x1A, 0x1A, 0x35) if r == 0 else
            RGBColor(0x12, 0x12, 0x28) if r % 2 == 0 else
            RGBColor(0x0F, 0x0F, 0x22)
        )
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        p.space_before = Pt(4)
        p.space_after  = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(13)
        run.font.bold = bold
        run.font.name = "Calibri"
        if color:
            run.font.color.rgb = color
        elif text == "✓":
            run.font.color.rgb = COL_GREEN
        elif text == "—":
            run.font.color.rgb = COL_MUTED
        else:
            run.font.color.rgb = COL_TEXT

    # Header row
    for c, h in enumerate(headers):
        _cell(0, c, h, bold=True, color=COL_MUTED, center=True)

    # Data rows
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            _cell(r, c, val, center=(c > 0))

    add_notes(slide, (
        ".cmd and .mac files get everything: highlighting, autocomplete, hover docs, "
        "linting, code navigation, and direct execution. Python gets the same. "
        "Solver files and output files get syntax highlighting only."
    ))
    return slide


def slide_19_section_adams_integration(prs):
    """Section: Adams Integration"""
    slide = prs.slides.add_slide(layout(prs, "Segue Slide Blue"))
    slide.placeholders[0].text = "Adams Integration"
    slide.placeholders[1].text = "Connect to a running Adams session"
    add_notes(slide, "You can connect to Adams View in several ways to run and debug code.")
    return slide


def slide_20_run_in_adams(prs):
    """Run in Adams View — GIF + keyboard shortcut."""
    slide = prs.slides.add_slide(layout(prs, "Title Only"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Run in Adams View"
    _style_title(slide)
    add_textbox(slide, "Run files or selected code.",
        x=Emu(446313), y=Emu(730000), cx=Emu(int(SLIDE_W * 0.8)),
        cy=Emu(300000), font_size=16, color=COL_MUTED)

    add_gif(slide, ASSETS_DIR / "run_selection_in_adams.gif",
        x=Emu(int(SLIDE_W * 0.05)), y=Emu(1200000),
        cx=Emu(int(SLIDE_W * 0.90)), cy=Emu(int(SLIDE_H * 0.66)))

    kbd_box = add_textbox(
        slide, "Ctrl+K  Ctrl+R  →  Executes directly in Adams View",
        x=Emu(int(SLIDE_W * 0.10)), y=Emu(int(SLIDE_H * 0.88)),
        cx=Emu(int(SLIDE_W * 0.80)), cy=Emu(350000),
        font_size=16, color=COL_ACCENT, align=PP_ALIGN.CENTER,
        font_name="Courier New",
    )
    hide_shape(kbd_box)
    add_appear_animations(slide, [kbd_box.shape_id])

    add_notes(slide, (
        "Select code in VS Code, press Ctrl+K Ctrl+R, and it executes directly "
        "in Adams View. No copy-paste. No switching windows. Your editor is your "
        "Adams console now. This works for both CMD and Python files."
    ))
    return slide


def slide_21_python_debug(prs):
    """Python Debugging in Adams — GIF."""
    slide = prs.slides.add_slide(layout(prs, "Title Only"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Python Debugging in Adams"
    _style_title(slide)
    add_textbox(
        slide,
        "Set breakpoints. Inspect variables. Step through code.  Running inside Adams.",
        x=Emu(446313), y=Emu(730000), cx=Emu(int(SLIDE_W * 0.8)),
        cy=Emu(300000), font_size=16, color=COL_MUTED)
    add_gif(slide, ASSETS_DIR / "debug_adams.gif",
        x=Emu(int(SLIDE_W * 0.05)), y=Emu(1200000),
        cx=Emu(int(SLIDE_W * 0.90)), cy=Emu(int(SLIDE_H * 0.73)))
    add_notes(slide, (
        "Full Python debugging inside Adams View. Set a breakpoint, click "
        "'Debug in Adams', and the debugger attaches to the running Adams process. "
        "When your script hits the breakpoint, execution pauses. You can inspect "
        "variables, step through code, evaluate expressions — the full debugging "
        "experience, inside Adams. No more print statements and prayer."
    ))
    return slide


def slide_22_section_whats_coming(prs):
    """Section: What's Coming"""
    slide = prs.slides.add_slide(layout(prs, "Segue Slide Blue"))
    slide.placeholders[0].text = "What's Coming"
    slide.placeholders[1].text = "AI-powered Adams scripting"
    add_notes(slide, "We've seen what the extension does today. Let me give you a glimpse of what's coming.")
    return slide


def slide_23_ai_agents(prs):
    """Teaching AI Agents to Use Adams — 2-column."""
    slide = prs.slides.add_slide(layout(prs, "2-Column Heavy Content"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Teaching AI Agents to Use Adams"
    _style_title(slide)

    # Sub-title row (idx=12)
    try:
        slide.placeholders[12].text = ""
    except KeyError:
        pass

    # Left column: MCP Servers (idx=13)
    # Right column: Agent Skills (idx=15)
    # We'll use free-form textboxes for better control over this complex slide.
    # Clear the placeholder text
    try:
        slide.placeholders[13].text = ""
        slide.placeholders[15].text = ""
    except KeyError:
        pass

    col_left_x  = Emu(446313)
    col_right_x = Emu(6095998)
    col_y       = Emu(1100000)
    col_w       = Emu(5528441)
    col_h       = Emu(5300000)

    # ---------- LEFT: MCP Servers ----------
    add_textbox(slide, "Bundled MCP Servers",
        x=col_left_x, y=col_y, cx=col_w, cy=Emu(300000),
        font_size=12, bold=True, color=COL_MUTED)

    def_y = col_y + Emu(320000)
    def_box = slide.shapes.add_textbox(col_left_x, def_y, col_w, Emu(700000))
    def_box.fill.solid()
    def_box.fill.fore_color.rgb = RGBColor(0x10, 0x10, 0x28)
    tf = def_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("MCP (Model Context Protocol) is an open standard that lets AI agents "
              "call external tools — like running Adams commands — from inside the "
              "chat window.")
    r.font.size = Pt(10)
    r.font.color.rgb = COL_MUTED
    r.font.name = "Calibri"

    mcp_servers = [
        ("Adams View", "live session",
         ["adams_run_cmd", "adams_run_python", "adams_load_file",
          "adams_evaluate_expression", "adams_export_model_cmd",
          "adams_create_simulation_script", "adams_submit_simulation",
          "adams_run_batch", "adams_batch_status", "adams_read_session_log",
          "adams_get_model_names", "adams_launch_view"]),
        ("Adams CMD Linter", "static analysis",
         ["adams_lint_cmd_text", "adams_lint_cmd_file", "adams_lookup_command"]),
    ]

    sy = def_y + Emu(750000)
    for srv_name, srv_tag, tools in mcp_servers:
        panel = slide.shapes.add_textbox(col_left_x, sy, col_w, Emu(1100000))
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x1E)
        tf = panel.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{srv_name}  [{srv_tag}]"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xA5, 0xB4, 0xFC)
        r.font.name = "Calibri"
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = "  ".join(tools)
        r2.font.size = Pt(9)
        r2.font.color.rgb = RGBColor(0x88, 0x99, 0xBB)
        r2.font.name = "Courier New"
        sy += Emu(1180000)

    # ---------- RIGHT: Agent Skills ----------
    add_textbox(slide, "Bundled Agent Skills",
        x=col_right_x, y=col_y, cx=col_w, cy=Emu(300000),
        font_size=12, bold=True, color=COL_MUTED)

    def_box2 = slide.shapes.add_textbox(col_right_x, def_y, col_w, Emu(700000))
    def_box2.fill.solid()
    def_box2.fill.fore_color.rgb = RGBColor(0x10, 0x10, 0x28)
    tf2 = def_box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = ("Agent skills are domain knowledge packs that teach Copilot "
               "Adams-specific concepts, patterns, syntax, and best practices.")
    r2.font.size = Pt(10)
    r2.font.color.rgb = COL_MUTED
    r2.font.name = "Calibri"

    skills = [
        ("adams-cmd-model-builder",    "Build models in Adams CMD syntax"),
        ("adams-python-model-builder", "Build models with the Adams Python API"),
        ("adams-flex",                 "Flexible bodies and MNF files"),
        ("adams-simulation-debugger",  "Diagnose convergence failures"),
        ("adams-subroutine-writer",    "Fortran & C user subroutines"),
    ]

    sy2 = def_y + Emu(750000)
    for skill_name, skill_desc in skills:
        row_box = slide.shapes.add_textbox(col_right_x, sy2, col_w, Emu(700000))
        row_box.fill.solid()
        row_box.fill.fore_color.rgb = RGBColor(0x0F, 0x0F, 0x26)
        tf_r = row_box.text_frame
        tf_r.word_wrap = True
        pr = tf_r.paragraphs[0]
        rr1 = pr.add_run()
        rr1.text = skill_name
        rr1.font.size = Pt(11)
        rr1.font.bold = True
        rr1.font.color.rgb = RGBColor(0xE0, 0xE7, 0xFF)
        rr1.font.name = "Courier New"
        pr2 = tf_r.add_paragraph()
        rr2 = pr2.add_run()
        rr2.text = skill_desc
        rr2.font.size = Pt(10)
        rr2.font.color.rgb = COL_MUTED
        rr2.font.name = "Calibri"
        sy2 += Emu(740000)

    add_notes(slide, (
        "The extension ships two MCP servers that let AI assistants like GitHub "
        "Copilot call into Adams directly. The Adams View server gives agents a "
        "live connection to a running Adams session. The CMD Linter server gives "
        "agents static analysis tools. On top of that, the extension bundles five "
        "agent skills — domain knowledge packs that teach Copilot how to think "
        "like an Adams engineer for specific tasks."
    ))
    return slide


def slide_24_copilot_example(prs):
    """Copilot Agent Example — chat placeholder + wind turbine image."""
    slide = prs.slides.add_slide(layout(prs, "Title Only"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Copilot Agent Example"
    _style_title(slide)

    left_x  = Emu(int(SLIDE_W * 0.04))
    right_x = Emu(int(SLIDE_W * 0.52))
    chat_y  = Emu(1150000)
    col_w   = Emu(int(SLIDE_W * 0.44))
    col_h   = Emu(int(SLIDE_H * 0.78))

    # --- Left: styled chat panel ---
    chat_bg = slide.shapes.add_textbox(left_x, chat_y, col_w, col_h)
    chat_bg.fill.solid()
    chat_bg.fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x28)
    tf = chat_bg.text_frame
    tf.word_wrap = True

    chat_lines = [
        ("You", False,
         "Build a wind turbine model in Adams. "
         "The tower is 80m tall, the rotor diameter is 100m. "
         "Add a revolute joint at the hub with gravity."),
        ("Copilot", True,
         "I'll build this step by step using the Adams CMD model-builder skill."),
        ("Copilot", True,
         "[Looks up PART CREATE syntax via adams_lookup_command]"),
        ("Copilot", True,
         "[Runs PART CREATE commands in Adams View via adams_run_cmd]"),
        ("Copilot", True,
         "Tower and nacelle parts created. "
         "Adding rotor with JOINT CREATE revolute..."),
        ("Copilot", True,
         "Model complete. Running a 10-second dynamic simulation."),
        ("You", False,
         "The simulation failed. Can you diagnose it?"),
        ("Copilot", True,
         "[Reads adams.msg via adams_read_session_log, identifies joint DOF error]"),
        ("Copilot", True,
         "Fixed the constraint. Re-running — simulation successful."),
    ]

    first = True
    for speaker, is_copilot, text in chat_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        r_speaker = p.add_run()
        r_speaker.text = f"{speaker}: "
        r_speaker.font.size = Pt(9)
        r_speaker.font.bold = True
        r_speaker.font.color.rgb = (COL_ACCENT if is_copilot else RGBColor(0x80, 0xD0, 0x80))
        r_speaker.font.name = "Calibri"
        r_text = p.add_run()
        r_text.text = text
        r_text.font.size = Pt(9)
        r_text.font.color.rgb = COL_TEXT
        r_text.font.name = "Calibri"
        # spacing
        p2 = tf.add_paragraph()
        p2.add_run().text = ""

    # --- Right: wind turbine model image ---
    add_image(slide, ASSETS_DIR / "wind_turbine_model.png",
        x=right_x, y=chat_y, cx=col_w, cy=col_h)

    add_notes(slide, (
        "Here's what it looks like in practice. You describe what you want in "
        "plain English. Copilot looks up the Adams command syntax, validates "
        "the expressions, and executes them in the running Adams session — "
        "all without you leaving the chat panel."
    ))
    return slide


def slide_25_section_get_started(prs):
    """Section: Get Started"""
    slide = prs.slides.add_slide(layout(prs, "Segue Slide Blue"))
    slide.placeholders[0].text = "Get Started"
    slide.placeholders[1].text = "It takes two minutes."
    return slide


def slide_26_installation(prs):
    """Installation — 2-column: steps + settings.json."""
    slide = prs.slides.add_slide(layout(prs, "2-Column Heavy Content"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Installation"
    _style_title(slide)
    try:
        slide.placeholders[12].text = ""
        slide.placeholders[13].text = ""
        slide.placeholders[15].text = ""
    except KeyError:
        pass

    col_left_x  = Emu(446313)
    col_right_x = Emu(6095998)
    col_y       = Emu(1150000)
    col_w       = Emu(5528441)

    # --- LEFT: installation steps ---
    add_textbox(slide, "From VS Code",
        x=col_left_x, y=col_y, cx=col_w, cy=Emu(350000),
        font_size=18, bold=True, color=COL_TEXT, font_name="Calibri Light")

    steps = [
        "1.  Open Extensions panel  (Ctrl+Shift+X)",
        '2.  Search  "MSC Adams"',
        "3.  Click  Install",
        "4.  Done.",
    ]
    sy = col_y + Emu(450000)
    for step in steps:
        add_textbox(slide, step,
            x=col_left_x + Emu(200000), y=sy, cx=col_w - Emu(200000), cy=Emu(500000),
            font_size=16, color=COL_TEXT)
        sy += Emu(520000)

    # Marketplace image
    mp_path = ASSETS_DIR / "vs_marketplace.png"
    if mp_path.exists():
        slide.shapes.add_picture(str(mp_path),
            col_left_x, sy + Emu(200000), height=Emu(800000))

    # --- RIGHT: settings.json code block ---
    add_textbox(slide, "Quick Setup",
        x=col_right_x, y=col_y, cx=col_w, cy=Emu(350000),
        font_size=18, bold=True, color=COL_TEXT, font_name="Calibri Light")

    settings_lines = [
        "// settings.json",
        "{",
        '  // Point to Adams installation',
        '  "msc-adams.adamsLaunchCommand":',
        r'    "C:\Program Files\MSC.Software\Adams\2024_2\...\mdi.bat",',
        "",
        "  // Enable the CMD linter",
        '  "msc-adams.linter.enabled": true,',
        "",
        "  // Scan workspace for macros",
        '  "msc-adams.linter.scanWorkspaceMacros": true',
        "}",
    ]
    add_code_block(slide, settings_lines,
        x=col_right_x, y=col_y + Emu(450000),
        cx=col_w, cy=Emu(4700000),
        font_size=11)

    add_notes(slide, (
        "Installation takes 30 seconds. Open the Extensions panel, search "
        "'MSC Adams', click Install. That's it. For the full experience, point "
        "the extension to your Adams installation and enable the linter and macro "
        "scanning. Three settings, and you're getting everything we just demonstrated."
    ))
    return slide


def slide_27_cta(prs):
    """Install it today — final CTA."""
    slide = prs.slides.add_slide(layout(prs, "Title with Tighter Content Leading"))
    set_bg(slide, "#0f0c29")
    slide.placeholders[0].text = "Install it today."
    _style_title(slide, font_size=40)

    add_textbox(slide, "Open your next .cmd file in VS Code.",
        x=Emu(446313), y=Emu(1050000), cx=Emu(int(SLIDE_W * 0.88)),
        cy=Emu(400000), font_size=20, color=COL_MUTED, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Let me know what breaks.",
        x=Emu(446313), y=Emu(1430000), cx=Emu(int(SLIDE_W * 0.88)),
        cy=Emu(380000), font_size=16, color=COL_MUTED, align=PP_ALIGN.CENTER)

    links = [
        ("VS Code Marketplace", "Search \"MSC Adams\"",
         "https://marketplace.visualstudio.com/items?itemName=savvyanalyst.msc-adams"),
        ("GitHub", "bthornton191/adams_vscode",
         "https://github.com/bthornton191/adams_vscode"),
        ("Feedback", "GitHub Issues",
         "https://github.com/bthornton191/adams_vscode/issues"),
        ("Contributions Welcome", "Open source",
         "https://github.com/bthornton191/adams_vscode/pulls"),
    ]

    total_w  = SLIDE_W - 2 * 446313
    col_w_ea = total_w // len(links)
    x_start  = 446313
    y_link   = int(SLIDE_H * 0.42)

    for label, sub, url in links:
        bx = slide.shapes.add_textbox(
            Emu(x_start), Emu(y_link), Emu(col_w_ea), Emu(int(SLIDE_H * 0.50)))
        tf = bx.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = label
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = COL_ACCENT
        r1.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(12)
        r2.font.color.rgb = COL_MUTED
        r2.font.name = "Calibri"
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = url
        r3.font.size = Pt(10)
        r3.font.color.rgb = RGBColor(0x80, 0xA0, 0xFF)
        r3.font.name = "Calibri"
        x_start += col_w_ea

    add_notes(slide, (
        "Install it today. Open your next .cmd file in VS Code. And let me know "
        "what breaks. Find it on the VS Code Marketplace — just search 'MSC Adams'. "
        "The source is on GitHub. If something doesn't work right, open a GitHub "
        "issue. And if you want to contribute, it's fully open source — jump in."
    ))
    return slide


# ---------------------------------------------------------------------------
# Internal: style the title placeholder consistently
# ---------------------------------------------------------------------------
def _style_title(slide, font_size: int = 28) -> None:
    try:
        ph = slide.placeholders[0]
    except (KeyError, IndexError):
        return
    for para in ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = COL_WHITE
            run.font.name = "Calibri Light"


# ---------------------------------------------------------------------------
# Patch content types to add mp4/gif/avi if missing
# ---------------------------------------------------------------------------
def _patch_content_types(pptx_path: Path) -> None:
    """Ensure [Content_Types].xml includes mp4, gif, avi extensions."""
    needed = {
        "mp4": 'Extension="mp4" ContentType="video/mp4"',
        "avi": 'Extension="avi" ContentType="video/avi"',
        "gif": 'Extension="gif" ContentType="image/gif"',
    }
    with zipfile.ZipFile(pptx_path, "r") as z:
        ct_data = z.read("[Content_Types].xml")
        all_files = z.infolist()
        all_data  = {item.filename: z.read(item.filename) for item in all_files}

    ct_str = ct_data.decode("utf-8")
    modified = False
    for ext, snippet in needed.items():
        if f'Extension="{ext}"' not in ct_str:
            ct_str = ct_str.replace(
                "</Types>",
                f'  <Default {snippet}/>\n</Types>'
            )
            modified = True

    if modified:
        all_data["[Content_Types].xml"] = ct_str.encode("utf-8")
        tmp = pptx_path.with_suffix(".tmp.pptx")
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as z_out:
            for filename, data in all_data.items():
                z_out.writestr(filename, data)
        tmp.replace(pptx_path)


# ===========================================================================
# Main
# ===========================================================================
def main():
    print("Loading template...")
    prs = load_template(POTX_PATH)

    print("Clearing template sample slides...")
    clear_slides(prs)

    print("Building slides...")
    builders = [
        (1,  "Title",                          slide_01_title),
        (2,  "Agenda",                         slide_02_agenda),
        (3,  "Section: The Editor Gap",        slide_03_section_editor_gap),
        (4,  "The Adams GUI Macro Editor",     slide_04_aview_editor),
        (5,  "The Notepad++ Upgrade",          slide_05_notepadpp),
        (6,  "What's Missing?",                slide_06_whats_missing),
        (7,  "Section: Code Editing",          slide_07_section_code_editing),
        (8,  "Syntax Highlighting",            slide_08_syntax_highlighting),
        (9,  "Autocomplete – CMD",             slide_09_autocomplete_cmd),
        (10, "Autocomplete – Python",          slide_10_autocomplete_python),
        (11, "Hover Docs – CMD",               slide_11_hover_cmd),
        (12, "Hover Docs – Python",            slide_12_hover_python),
        (13, "Linting – Demo",                 slide_13_linting_demo),
        (14, "Linting – CMD",                  slide_14_linting_cmd),
        (15, "Linting – Python",               slide_15_linting_python),
        (16, "Code Navigation",                slide_16_code_nav),
        (17, "Custom Macros",                  slide_17_custom_macros),
        (18, "Supported File Types",           slide_18_file_types),
        (19, "Section: Adams Integration",     slide_19_section_adams_integration),
        (20, "Run in Adams View",              slide_20_run_in_adams),
        (21, "Python Debugging",               slide_21_python_debug),
        (22, "Section: What's Coming",         slide_22_section_whats_coming),
        (23, "Teaching AI Agents",             slide_23_ai_agents),
        (24, "Copilot Agent Example",          slide_24_copilot_example),
        (25, "Section: Get Started",           slide_25_section_get_started),
        (26, "Installation",                   slide_26_installation),
        (27, "Install it today",               slide_27_cta),
    ]

    for num, name, fn in builders:
        print(f"  [{num:02d}/27] {name}")
        fn(prs)

    print(f"Saving to {OUTPUT_PATH}...")
    prs.save(str(OUTPUT_PATH))

    print("Patching content types for video/gif...")
    _patch_content_types(OUTPUT_PATH)

    print(f"\nDone! -> {OUTPUT_PATH}")
    print(f"File size: {OUTPUT_PATH.stat().st_size / 1024 / 1024:.1f} MB")


if __name__ == "__main__":
    main()
